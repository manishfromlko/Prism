"""
Two slides:
  1. Full Load   — Breadth-First Airflow DAG (4 stage columns + proportional timeline)
  2. Incremental — Depth-First Per-Workspace (trigger → hash-guard → loop + comparison)

Numbers derived from actual run on 2026-05-03:
  211 artifacts ingested  : ~2 s
  190 artifacts embedded  : ~11 s (6 batches of 32, ~1.83 s/batch)
  211 summaries generated : 9 min 36 s  (22.0/min rate)
  5 profiles generated    : ~20 s       (0.25 ws/s rate)

Run:  python generate_load_strategy_slides.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Load_Strategy_Full_and_Incremental.pptx"

# ── Palette ────────────────────────────────────────────────────────────────────
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(248, 250, 252)
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(100, 116, 139)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(180, 100, 0)
RED    = RGBColor(185, 28,  28)
PINK   = RGBColor(190, 24,  93)
PURPLE = RGBColor(109, 40,  217)
TEAL   = RGBColor(13,  148, 136)
SLATE  = RGBColor(148, 163, 184)
DARKSL = RGBColor(71,  85,  105)
PANEL  = RGBColor(241, 245, 249)
BLUEBG = RGBColor(219, 234, 254)
GREENBG= RGBColor(220, 252, 231)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)
TEALB  = RGBColor(204, 251, 241)
WARN   = RGBColor(254, 226, 226)
REDBRD = RGBColor(220, 38,  38)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _shape(slide, kind, x, y, w, h):
    return slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))


def box(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=0.8,
        dash=False, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = _shape(slide, kind, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    if dash:
        s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius:
        s.adjustments[0] = 0.04
    return s


def tb(slide, x, y, w, h, text, size=10, color=INK, bold=False,
       align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color
    return b


def multi(slide, x, y, w, h, rows):
    """rows = [(text, size, color, bold, align), ...]"""
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold, align) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.alignment = align
        p.space_after = Pt(1)
    return b


def arr(slide, x1, y1, x2, y2, color=DARKSL, lw=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(lw)
    return c


def dot(slide, x, y, r, color=BLUE):
    s = _shape(slide, MSO_SHAPE.OVAL, x, y, r, r)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def bullet_rows(slide, x, y, w, items, size=8.8, icolor=DARKSL, dot_color=MUTED):
    yy = y
    for item in items:
        dot(slide, x + 0.03, yy + 0.04, 0.08, dot_color)
        tb(slide, x + 0.16, yy, w - 0.18, 0.25, item, size, icolor)
        yy += 0.25
    return yy


def brand(slide, label="Kubeflow Workspace Intelligence LLMOps"):
    b = box(slide, 0, 0, 13.33, 0.30, INK, INK)
    b.line.fill.background()
    for i, c in enumerate([BLUE, GREEN, AMBER, PINK]):
        a = box(slide, i * 3.3325, 0.30, 3.3325, 0.06, c, c)
        a.line.fill.background()
    tb(slide, 0.5, 7.14, 12.5, 0.24, label,
       size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def pill(slide, x, y, w, h, text, fill, line, tcolor=WHITE, tsize=9.5, tbold=True):
    s = box(slide, x, y, w, h, fill, line, 1.0, radius=True)
    tb(slide, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
       text, tsize, tcolor, tbold, PP_ALIGN.CENTER)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — FULL LOAD
# ══════════════════════════════════════════════════════════════════════════════

STAGE_DATA = [
    # (num, short_name, color, bg, airflow, ops, src, dst, seq, par, api_stat, is_bn)
    dict(
        num="01", name="Data Ingestion", color=BLUE, bg=BLUEBG,
        airflow="500 workspace tasks  ·  20 concurrent  →  25 rounds",
        ops=[
            "Scan workspace dirs for .ipynb / .py / .sql / .md",
            "Extract text + notebook cell content",
            "MD5 content-hash fingerprinting",
            "Write ingestion_catalog.json (100K records)",
        ],
        src="Workspace files\n(dataset/  —  500 WS)",
        dst="ingestion_catalog.json",
        dst_color=BLUE,
        seq="~16 min  (sequential)",
        par="~1.5 min  (20 workers)",
        api_stat="No API calls  —  pure I/O",
        is_bn=False,
    ),
    dict(
        num="02", name="Artifact Embedding", color=PURPLE, bg=PURPBG,
        airflow="500 workspace tasks  ·  20 concurrent  →  25 rounds",
        ops=[
            "Load catalog  →  191 docs / workspace avg",
            "Chunk + guardrail filter (1 rejected / WS avg)",
            "OpenAI embed: batch_size=32  →  3,125 batches",
            "Upsert 100K vectors into Milvus (1536-dim)",
        ],
        src="ingestion_catalog.json\n(100K artifact records)",
        dst="kubeflow_artifacts\n(Milvus — 100K vectors)",
        dst_color=PURPLE,
        seq="~1 hr 35 min  (sequential)",
        par="~9 min  (20 workers)",
        api_stat="3,125 OpenAI embed calls  ·  ~1.83 s/batch",
        is_bn=False,
    ),
    dict(
        num="03", name="Artifact Summaries", color=REDBRD, bg=WARN,
        airflow="500 workspace tasks  ·  20 concurrent  →  25 rounds",
        ops=[
            "1 LLM call per artifact  (gpt-4o-mini)",
            "Max 220 output tokens  ·  JSON response",
            "Rate: 22 artifacts/min per worker",
            "Embed summaries  →  upsert to Milvus",
        ],
        src="ingestion_catalog.json\n(100K artifacts)",
        dst="artifact_summaries\n(Milvus — 100K entries)",
        dst_color=REDBRD,
        seq="~75 hrs 50 min  (sequential)",
        par="~4 hrs 2 min  (20 workers)",
        api_stat="100,000 LLM calls  ·  2.5 s avg latency",
        is_bn=True,
    ),
    dict(
        num="04", name="User Profiles", color=PINK, bg=PINKBG,
        airflow="500 workspace tasks  ·  20 concurrent  →  25 rounds",
        ops=[
            "Fetch artifact summaries per workspace",
            "1 LLM call per workspace  (gpt-4o-mini)",
            "Profile text + skill tags generated",
            "Embed + upsert into user_profiles",
        ],
        src="artifact_summaries\n(Milvus — per workspace)",
        dst="user_profiles\n(Milvus — 500 profiles)",
        dst_color=PINK,
        seq="~33 min  (sequential)",
        par="~4 min  (20 workers)",
        api_stat="500 LLM calls  ·  4 s avg latency",
        is_bn=False,
    ),
]

# Timeline proportional widths (minimum floor for visual clarity)
_times  = [1.5, 9, 242, 4]      # minutes with 20 workers
_total  = sum(_times)
_BAR_W  = 12.85
_floors = [0.30, 0.55, 0.0, 0.40]   # minimum widths for S1, S2, blank, S4
_S3_w   = _BAR_W - sum(_floors)     # S3 gets the rest → 11.6"
_BAR_WIDTHS = [_floors[0], _floors[1], _S3_w, _floors[3]]

def make_full_load(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    brand(slide)

    # ── Title + context ──────────────────────────────────────────────────────
    tb(slide, 0.25, 0.38, 13.0, 0.44,
       "Full Load  —  Breadth-First Airflow DAG",
       size=22, color=INK, bold=True)

    # Context row
    ctx_items = [
        ("500 Workspaces", BLUEBG, BLUE),
        ("×  200 Artifacts avg", AMBERBG, AMBER),
        ("=  100K Total", GREENBG, GREEN),
        ("20 Airflow Workers", PURPBG, PURPLE),
        ("Breadth-First: all WS → Stage N before Stage N+1", PANEL, DARKSL),
    ]
    cx = 0.25
    for label, fill, line in ctx_items:
        w = 1.35 if "Breadth" not in label else 3.4
        pill(slide, cx, 0.85, w, 0.30, label, fill, line, INK, 8, True)
        cx += w + 0.12

    # ── Stage columns ─────────────────────────────────────────────────────────
    SW = 3.05    # stage width
    AW = 0.18    # arrow width
    SX = [0.20, 0.20 + SW + AW,
          0.20 + 2*(SW + AW), 0.20 + 3*(SW + AW)]
    SY = 1.22
    SH = 4.22

    for idx, sd in enumerate(STAGE_DATA):
        x = SX[idx]
        c, bg = sd["color"], sd["bg"]

        # Outer box
        box(slide, x, SY, SW, SH, WHITE, c, 1.3, radius=True)

        # ── Header band ──────────────────────────────────────────────────────
        hdr = box(slide, x, SY, SW, 0.48, c, c, 0)
        hdr.line.fill.background()
        # Badge
        bdg = box(slide, x + 0.10, SY + 0.09, 0.32, 0.30, WHITE, WHITE, 0)
        bdg.line.fill.background()
        tb(slide, x + 0.10, SY + 0.09, 0.32, 0.30,
           sd["num"], 9, c, True, PP_ALIGN.CENTER)
        tb(slide, x + 0.48, SY + 0.11, SW - 0.58, 0.30,
           sd["name"], 12.5, WHITE, True)

        # ── Airflow tasks ─────────────────────────────────────────────────────
        airflow_y = SY + 0.54
        af_box = box(slide, x + 0.1, airflow_y, SW - 0.2, 0.45, PANEL, SLATE, 0.6)
        multi(slide, x + 0.14, airflow_y + 0.02, SW - 0.28, 0.42, [
            ("✦ Airflow", 7.5, DARKSL, True, PP_ALIGN.LEFT),
            (sd["airflow"], 8, INK, False, PP_ALIGN.LEFT),
        ])

        # ── Operations ────────────────────────────────────────────────────────
        ops_y = airflow_y + 0.52
        tb(slide, x + 0.12, ops_y, SW - 0.22, 0.22,
           "What happens:", 8, MUTED, True)
        yy = ops_y + 0.23
        for op in sd["ops"]:
            dot(slide, x + 0.14, yy + 0.045, 0.07, c)
            tb(slide, x + 0.27, yy, SW - 0.38, 0.24, op, 8.5, INK)
            yy += 0.25

        # ── Source → Milvus ───────────────────────────────────────────────────
        io_y = ops_y + 1.32
        # Source box
        src_b = box(slide, x + 0.1, io_y, SW - 0.2, 0.48, PANEL, SLATE, 0.7)
        multi(slide, x + 0.14, io_y + 0.03, SW - 0.28, 0.43, [
            ("IN: ", 8, MUTED, True, PP_ALIGN.LEFT),
            (sd["src"], 8, INK, False, PP_ALIGN.LEFT),
        ])
        # Arrow down
        arr(slide, x + SW/2, io_y + 0.48, x + SW/2, io_y + 0.62, c, 1.2)
        # Destination Milvus box
        dst_b = box(slide, x + 0.1, io_y + 0.62, SW - 0.2, 0.48,
                    bg, sd["dst_color"], 1.0)
        multi(slide, x + 0.14, io_y + 0.64, SW - 0.28, 0.43, [
            ("OUT → Milvus: ", 8, sd["dst_color"], True, PP_ALIGN.LEFT),
            (sd["dst"].split("\n")[-1], 8.5, INK, True, PP_ALIGN.LEFT),
        ])

        # ── Sequential time (muted) ───────────────────────────────────────────
        seq_y = io_y + 1.18
        tb(slide, x + 0.12, seq_y, SW - 0.22, 0.26,
           sd["seq"], 8.5, MUTED, False, PP_ALIGN.LEFT, italic=True)

        # ── Parallel time (highlighted) ───────────────────────────────────────
        par_y = seq_y + 0.28
        par_fill = WARN if sd["is_bn"] else GREENBG
        par_line = REDBRD if sd["is_bn"] else GREEN
        par_col  = REDBRD if sd["is_bn"] else GREEN
        pb = box(slide, x + 0.10, par_y, SW - 0.2, 0.44,
                 par_fill, par_line, 1.2, radius=True)
        multi(slide, x + 0.14, par_y + 0.03, SW - 0.28, 0.40, [
            ("⚡  20 workers:", 8.5, MUTED, True, PP_ALIGN.LEFT),
            (sd["par"].split("(")[0].strip(), 11.5, par_col, True, PP_ALIGN.LEFT),
        ])

        # ── API stat ──────────────────────────────────────────────────────────
        api_y = par_y + 0.52
        tb(slide, x + 0.12, api_y, SW - 0.22, 0.26,
           sd["api_stat"], 8, MUTED, False)

        # ── Bottleneck badge ──────────────────────────────────────────────────
        if sd["is_bn"]:
            bn_y = api_y + 0.28
            bn = box(slide, x + 0.10, bn_y, SW - 0.2, 0.26,
                     WARN, REDBRD, 1.0, radius=True)
            tb(slide, x + 0.14, bn_y + 0.02, SW - 0.28, 0.24,
               "⚠  CRITICAL PATH — drives total runtime",
               8.5, REDBRD, True, PP_ALIGN.CENTER)

        # ── Arrow to next stage ───────────────────────────────────────────────
        if idx < 3:
            ax = x + SW
            mid_y = SY + SH / 2
            arr(slide, ax, mid_y, ax + AW, mid_y, c, 2.2)
            # arrowhead
            tri = _shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE,
                         ax + AW - 0.10, mid_y - 0.07, 0.12, 0.14)
            tri.fill.solid(); tri.fill.fore_color.rgb = c
            tri.line.fill.background(); tri.rotation = 90

    # ══════════════════════════════════════════════════════════════════════════
    # PROPORTIONAL TIMELINE BAR
    # ══════════════════════════════════════════════════════════════════════════
    BAR_Y = SY + SH + 0.18
    bar_colors = [BLUE, PURPLE, REDBRD, PINK]
    bar_labels = ["~1.5 m", "~9 m", "~4 hrs 2 min  (bottleneck)", "~4 m"]
    bx = 0.24

    for i, (w, col, label) in enumerate(zip(_BAR_WIDTHS, bar_colors, bar_labels)):
        b = box(slide, bx, BAR_Y, w, 0.36, col, col)
        b.line.fill.background()
        # label inside or below
        if w > 0.8:
            tb(slide, bx + 0.08, BAR_Y + 0.07, w - 0.12, 0.24,
               label, 9.5, WHITE, True, PP_ALIGN.CENTER)
        else:
            tb(slide, bx, BAR_Y + 0.38, w + 0.1, 0.22,
               label, 7.5, col, True, PP_ALIGN.CENTER)
        bx += w

    tb(slide, 0.24, BAR_Y - 0.22, 6.0, 0.20,
       "Proportional runtime  (20 Airflow workers):", 8, MUTED, False)
    tb(slide, 9.5, BAR_Y - 0.22, 3.7, 0.20,
       "← Stage 3 = 94.5% of total runtime",
       8, REDBRD, True, PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════════════════════
    # TOTAL + KEY METRICS
    # ══════════════════════════════════════════════════════════════════════════
    MET_Y = BAR_Y + 0.64

    pill(slide, 0.24, MET_Y, 2.0, 0.34, "Total: ~4 hrs 16 min",
         GREENBG, GREEN, INK, 9.5, True)

    metrics = [
        ("100,000 LLM calls", PURPBG, PURPLE),
        ("3,125 embed batches", BLUEBG, BLUE),
        ("4 Milvus collections updated", AMBERBG, AMBER),
        ("Bottleneck: Stage 3 — Summaries", WARN, REDBRD),
    ]
    mx = 2.38
    for label, fill, line in metrics:
        w = 2.35 if "Bottleneck" in label else 2.1
        pill(slide, mx, MET_Y, w, 0.34, label, fill, line, INK, 8.5, False)
        mx += w + 0.10

    # ── Assumptions ───────────────────────────────────────────────────────────
    tb(slide, 0.24, MET_Y + 0.42, 13.0, 0.22,
       "Basis: observed 2026-05-03 run — 211 artifacts/2 s (ingest), "
       "190 artifacts/11 s (embed, 6×32-batch), "
       "211 summaries/9m 36s (22/min).  "
       "Extrapolated to 100K.  "
       "OpenAI Tier 4: gpt-4o-mini 10K RPM · embeddings 5K RPM · "
       "20 Airflow workers · batch_size=32.",
       size=7.5, color=MUTED, italic=True)

    slide.notes_slide.notes_text_frame.text = (
        "FULL LOAD — Breadth-First. All 500 workspaces complete each stage before the next starts.\n"
        "Stage 1: 1.5 min  Stage 2: 9 min  Stage 3: 4h 2min (BOTTLENECK)  Stage 4: 4 min\n"
        "Total: ~4 hrs 16 min with 20 Airflow workers.\n"
        "Numbers derived from actual run on 2026-05-03 then extrapolated to 100K artifacts."
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INCREMENTAL LOAD
# ══════════════════════════════════════════════════════════════════════════════

def make_incremental(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    brand(slide)

    # ── Title + scenario context ──────────────────────────────────────────────
    tb(slide, 0.25, 0.38, 9.5, 0.44,
       "Incremental Load  —  Depth-First Per-Workspace",
       size=22, color=INK, bold=True)

    scenario_pills = [
        ("Daily trigger", PANEL, DARKSL),
        ("5% WS changed → 25 workspaces", BLUEBG, BLUE),
        ("10% artifacts/WS → 500 artifacts", AMBERBG, AMBER),
        ("5 Airflow workers", PURPBG, PURPLE),
    ]
    cx = 0.25
    for label, fill, line in scenario_pills:
        w = 1.1 if len(label) < 18 else 2.45
        pill(slide, cx, 0.85, w, 0.30, label, fill, line, INK, 8, True)
        cx += w + 0.10

    # ════════════════════════════════════════════════
    # LEFT FLOW (x=0.25 to 8.8)
    # ════════════════════════════════════════════════
    LW = 8.3    # left column width
    FX = 0.25   # flow left edge

    # ── TRIGGER ──────────────────────────────────────────────────────────────
    TRG_Y = 1.22
    trg = box(slide, FX + 2.55, TRG_Y, 3.2, 0.54, GREENBG, GREEN, 1.3, radius=True)
    multi(slide, FX + 2.63, TRG_Y + 0.04, 3.04, 0.48, [
        ("Airflow Sensor  —  Daily Schedule", 10.5, GREEN, True, PP_ALIGN.CENTER),
        ("or  file-system change event", 8.5, MUTED, False, PP_ALIGN.CENTER),
    ])

    arr(slide, FX + 4.15, TRG_Y + 0.54, FX + 4.15, TRG_Y + 0.78, GREEN, 1.5)
    tb(slide, FX + 4.25, TRG_Y + 0.56, 1.5, 0.22,
       "triggers", 8, MUTED, italic=True)

    # ── STAGE 0: CHANGE DETECTION ─────────────────────────────────────────────
    DET_Y = TRG_Y + 0.78
    det = box(slide, FX + 0.5, DET_Y, 7.3, 0.72, BLUEBG, BLUE, 1.3, radius=True)
    multi(slide, FX + 0.62, DET_Y + 0.06, 7.1, 0.64, [
        ("Stage 0 — Change Detection  (~40 sec)", 10.5, BLUE, True, PP_ALIGN.LEFT),
        ("Scan 25 changed workspaces  ·  "
         "Compare MD5 content hashes  ·  "
         "Identify 500 changed artifacts  ·  "
         "180 unchanged per WS → SKIP",
         8.5, INK, False, PP_ALIGN.LEFT),
    ])

    # Split arrow: 25 WS with changes → per-workspace loop
    arr(slide, FX + 4.15, DET_Y + 0.72, FX + 4.15, DET_Y + 0.98, BLUE, 1.5)

    # ── PER-WORKSPACE LOOP (dashed boundary) ──────────────────────────────────
    LOOP_Y = DET_Y + 0.98
    LOOP_H = 3.82
    loop_box = box(slide, FX + 0.15, LOOP_Y, LW, LOOP_H,
                   fill=None, line=DARKSL, lw=1.4, dash=True, radius=True)

    tb(slide, FX + 0.25, LOOP_Y + 0.06, 4.5, 0.24,
       "Per-Workspace Loop  —  25 workspaces  ·  5 concurrent  →  5 rounds",
       size=9, color=DARKSL, bold=True)

    # Hash guard decision
    HASH_Y = LOOP_Y + 0.38
    hg = box(slide, FX + 0.35, HASH_Y, 2.45, 0.64, PANEL, SLATE, 1.0, radius=True)
    multi(slide, FX + 0.43, HASH_Y + 0.06, 2.3, 0.54, [
        ("Hash Guard", 9.5, DARKSL, True, PP_ALIGN.CENTER),
        ("Per artifact: compare MD5", 8, MUTED, False, PP_ALIGN.CENTER),
    ])

    # Skip unchanged
    skip = box(slide, FX + 0.35, HASH_Y + 0.76, 2.45, 0.50,
               GREENBG, GREEN, 0.9, radius=True)
    multi(slide, FX + 0.43, HASH_Y + 0.80, 2.3, 0.42, [
        ("SKIP  ✓  (unchanged)", 9, GREEN, True, PP_ALIGN.CENTER),
        ("~180 artifacts / WS = 90%", 8, MUTED, False, PP_ALIGN.CENTER),
    ])
    arr(slide, FX + 1.58, HASH_Y + 0.64, FX + 1.58, HASH_Y + 0.76, GREEN, 1.2)

    # Changed artifacts branch
    arr(slide, FX + 2.80, HASH_Y + 0.32, FX + 3.15, HASH_Y + 0.32, REDBRD, 1.3)
    tb(slide, FX + 2.82, HASH_Y + 0.08, 0.38, 0.22,
       "changed", 7.5, REDBRD, italic=True)

    # ── 3 processing stages in the loop ───────────────────────────────────────
    PROC_X  = FX + 3.18
    PROC_W  = 4.9
    PROC_H  = 0.60
    GAP     = 0.28

    proc_stages = [
        ("Stage 2 — Re-Embed",
         "500 artifacts  →  16 API batches (batch_size=32)  →  Milvus upsert",
         "~1 min", PURPLE, PURPBG, "kubeflow_artifacts  (partial upsert)"),
        ("Stage 3 — Re-Summarise",
         "500 LLM calls (gpt-4o-mini)  ·  5 workers × 100 artifacts  ·  ~50 s/round",
         "~5 min", REDBRD, WARN,   "artifact_summaries  (partial upsert)"),
        ("Stage 4 — Re-Profile",
         "25 workspace profiles regenerated from updated summaries",
         "~20 sec", PINK, PINKBG, "user_profiles  (25 upserted)"),
    ]

    py = HASH_Y
    for i, (title, desc, timing, col, bg, milvus) in enumerate(proc_stages):
        pb = box(slide, PROC_X, py, PROC_W, PROC_H + 0.22, bg, col, 1.1, radius=True)

        multi(slide, PROC_X + 0.12, py + 0.05, PROC_W - 0.2, PROC_H + 0.14, [
            (title, 10, col, True, PP_ALIGN.LEFT),
            (desc,  8.2, INK, False, PP_ALIGN.LEFT),
            (f"↳  {milvus}", 8, MUTED, False, PP_ALIGN.LEFT),
        ])

        # Timing badge
        pill(slide, PROC_X + PROC_W - 1.22, py + 0.53, 1.18, 0.26,
             timing, GREENBG, GREEN, INK, 9, True)

        if i < 2:
            arr(slide, PROC_X + PROC_W / 2, py + PROC_H + 0.22,
                PROC_X + PROC_W / 2, py + PROC_H + 0.22 + GAP + 0.02,
                col, 1.3)
        py += PROC_H + 0.22 + GAP

    # Total loop time badge
    loop_total_y = LOOP_Y + LOOP_H - 0.46
    pill(slide, FX + 0.3, loop_total_y, 2.6, 0.36,
         "Loop total (per 25 WS):  ~7 min",
         GREENBG, GREEN, INK, 9.5, True)
    tb(slide, FX + 3.1, loop_total_y + 0.06, 5.1, 0.26,
       "vs Full Load Stage 2–4: ~4 hrs 15 min  →  36× faster for 5% delta",
       9.5, GREEN, True)

    # ── DONE box ──────────────────────────────────────────────────────────────
    DONE_Y = LOOP_Y + LOOP_H + 0.14
    arr(slide, FX + 4.15, LOOP_Y + LOOP_H, FX + 4.15, DONE_Y, GREEN, 1.5)
    done = box(slide, FX + 2.55, DONE_Y, 3.2, 0.44, GREENBG, GREEN, 1.3, radius=True)
    tb(slide, FX + 2.63, DONE_Y + 0.08, 3.04, 0.30,
       "Done  —  Notify / log metadata",
       10, GREEN, True, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════
    # RIGHT COLUMN — COMPARISON TABLE (x=8.75 to 13.1)
    # ════════════════════════════════════════════════
    TX = 8.72
    TW = 4.42
    T_Y = 1.22

    # Table header
    th = box(slide, TX, T_Y, TW, 0.40, INK, INK)
    th.line.fill.background()
    tb(slide, TX + 0.1, T_Y + 0.07, TW - 0.2, 0.28,
       "Full Load  vs  Incremental", 11, WHITE, True, PP_ALIGN.CENTER)

    rows = [
        ("",               "Full Load",          "Incremental",     INK,    INK),
        ("Trigger",        "Weekly / Monthly",    "Daily / On-change", MUTED, MUTED),
        ("Scope",          "100K artifacts",      "~500 artifacts",  INK,    INK),
        ("Workspaces",     "500 (all)",           "~25  (5% delta)", MUTED,  MUTED),
        ("Artifacts",      "100,000",             "~500",            INK,    INK),
        ("Hash guard",     "No  (full rebuild)",  "Yes  (skip 90%)", MUTED,  GREEN),
        ("Workers",        "20",                  "5",               INK,    INK),
        ("Stage 1",        "~1.5 min",            "N/A  (unchanged)",MUTED,  GREEN),
        ("Stage 2 Embed",  "~9 min",              "~1 min",          BLUE,   GREEN),
        ("Stage 3 Summ.",  "~4 hrs 2 min ⚠",     "~5 min",          REDBRD, GREEN),
        ("Stage 4 Prof.",  "~4 min",              "~20 sec",         PINK,   GREEN),
        ("TOTAL",          "~4 hrs 16 min",       "~7 min",          REDBRD, GREEN),
    ]

    row_h = 0.34
    ry = T_Y + 0.40
    for i, (label, full, inc, fc, ic) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else PANEL
        if label == "TOTAL":
            bg = GREENBG
        row_box = box(slide, TX, ry, TW, row_h, bg, SLATE, 0.5)
        if label == "":
            # Sub-header row
            box(slide, TX, ry, TW, row_h, PANEL, SLATE, 0.5)
            tb(slide, TX + 0.08, ry + 0.06, 1.3, row_h - 0.1, "", 8, MUTED, True)
            tb(slide, TX + 1.48, ry + 0.06, 1.45, row_h - 0.1, full,  8.5, DARKSL, True, PP_ALIGN.CENTER)
            tb(slide, TX + 2.98, ry + 0.06, 1.40, row_h - 0.1, inc,   8.5, DARKSL, True, PP_ALIGN.CENTER)
        else:
            tb(slide, TX + 0.08, ry + 0.07, 1.35, row_h - 0.1, label, 8.5, DARKSL, True)
            tb(slide, TX + 1.48, ry + 0.07, 1.45, row_h - 0.1, full,  8.5, fc, label == "TOTAL", PP_ALIGN.CENTER)
            tb(slide, TX + 2.98, ry + 0.07, 1.40, row_h - 0.1, inc,   8.5, ic, label == "TOTAL", PP_ALIGN.CENTER)
            # vertical dividers
            arr(slide, TX + 1.45, ry, TX + 1.45, ry + row_h, SLATE, 0.5)
            arr(slide, TX + 2.95, ry, TX + 2.95, ry + row_h, SLATE, 0.5)
        ry += row_h

    # Column header divider
    arr(slide, TX, T_Y + 0.40, TX + TW, T_Y + 0.40, SLATE, 0.8)

    # Speedup callout
    sp_y = ry + 0.12
    sp = box(slide, TX, sp_y, TW, 0.52, GREENBG, GREEN, 1.3, radius=True)
    multi(slide, TX + 0.12, sp_y + 0.06, TW - 0.2, 0.42, [
        ("~36× faster for 5% daily delta", 12, GREEN, True, PP_ALIGN.CENTER),
        ("DFS atomicity: per-WS consistent state", 8.5, DARKSL, False, PP_ALIGN.CENTER),
    ])

    slide.notes_slide.notes_text_frame.text = (
        "INCREMENTAL LOAD — Depth-First. Each changed workspace flows through all stages before the next.\n"
        "Scenario: daily run, 5% workspaces changed = 25 WS, 10% artifacts per WS changed = 500 artifacts.\n"
        "Hash guard skips 90% of artifacts (unchanged).\n"
        "Total: ~7 min with 5 workers — vs ~4h 16m for full load = 36× faster.\n"
        "Key advantage: per-workspace atomicity. If workspace 15 fails, workspaces 1-14 are already committed."
    )


# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    make_full_load(prs)
    make_incremental(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
