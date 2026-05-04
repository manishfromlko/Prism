"""
Three slides:
  1. Retrieval Process   — intent classification → query rewrite → vector retrieval paths
  2. Generation Process  — context assembly → LLM generate → response format + tracing
  3. Evaluation Process  — Layer 1 heuristics / Layer 2 RAGAS / LLM judge / user feedback

Run:  python generate_eval_slides.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Retrieval_Generation_Evaluation.pptx"

# ── Palette ────────────────────────────────────────────────────────────────────
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(248, 250, 252)
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(100, 116, 139)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(180, 100, 0)
RED    = RGBColor(185, 28,  28)
PINK   = RGBColor(190, 24,  93)
PURPLE = RGBColor(109, 40,  217)
TEAL   = RGBColor(13,  148, 136)
SLATE  = RGBColor(148, 163, 184)
DARKSL = RGBColor(71,  85,  105)
PANEL  = RGBColor(241, 245, 249)
BLUEBG = RGBColor(219, 234, 254)
GREENBG= RGBColor(220, 252, 231)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)
TEALB  = RGBColor(204, 251, 241)
REDBRD = RGBColor(220, 38,  38)
ORANGEBG = RGBColor(255, 237, 213)
ORANGE = RGBColor(194, 65,  12)


# ── Helpers ────────────────────────────────────────────────────────────────────

def _shape(slide, kind, x, y, w, h):
    return slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))


def box(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=0.8, dash=False, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = _shape(slide, kind, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    if dash:
        s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius:
        s.adjustments[0] = 0.04
    return s


def tb(slide, x, y, w, h, text, size=10, color=INK, bold=False,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color
    return b


def multi(slide, x, y, w, h, rows):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold, align) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.alignment = align
        p.space_after = Pt(1)
    return b


def arr(slide, x1, y1, x2, y2, color=DARKSL, lw=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(lw)
    return c


def dot(slide, x, y, r, color=BLUE):
    s = _shape(slide, MSO_SHAPE.OVAL, x, y, r, r)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def pill(slide, x, y, w, h, text, fill, line, tcolor=WHITE, tsize=9, tbold=True):
    s = box(slide, x, y, w, h, fill, line, 1.0, radius=True)
    tb(slide, x + 0.05, y + 0.03, w - 0.1, h - 0.06,
       text, tsize, tcolor, tbold, PP_ALIGN.CENTER)
    return s


def brand(slide, label="Kubeflow Workspace Intelligence LLMOps"):
    b = box(slide, 0, 0, 13.33, 0.30, INK, INK)
    b.line.fill.background()
    for i, c in enumerate([BLUE, GREEN, AMBER, PINK]):
        a = box(slide, i * 3.3325, 0.30, 3.3325, 0.06, c, c)
        a.line.fill.background()
    tb(slide, 0.5, 7.14, 12.5, 0.24, label,
       size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def section_header(slide, x, y, w, h, title, fill=BLUE, tcolor=WHITE):
    box(slide, x, y, w, h, fill, fill, radius=True)
    tb(slide, x + 0.1, y + 0.04, w - 0.2, h - 0.08,
       title, 9.5, tcolor, True, PP_ALIGN.CENTER)


def step_box(slide, x, y, w, h, num, title, subtitle,
             num_fill=BLUE, bg=PANEL, line=SLATE):
    box(slide, x, y, w, h, bg, line, 0.8, radius=True)
    dot(slide, x + 0.1, y + 0.1, 0.28, num_fill)
    tb(slide, x + 0.12, y + 0.11, 0.24, 0.28,
       str(num), 10, WHITE, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.44, y + 0.07, w - 0.54, 0.22, title, 9, INK, True)
    tb(slide, x + 0.44, y + 0.26, w - 0.54, 0.30, subtitle, 7.5, MUTED, False)


# ── Slide 1: Retrieval Process ─────────────────────────────────────────────────

def slide_retrieval(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand(slide)

    # Title
    box(slide, 0.38, 0.42, 12.57, 0.46, BLUEBG, BLUE, 0)
    tb(slide, 0.5, 0.45, 12.3, 0.40,
       "Retrieval Process — Intent Classification → Query Rewrite → Vector Search",
       12.5, INK, True, PP_ALIGN.LEFT)

    # ── Column headers ─────────────────────────────────────────────────────────
    headers = [
        (0.38, "① User Query", BLUE,   BLUEBG),
        (2.90, "② Classify",   PURPLE, PURPBG),
        (5.42, "③ Rewrite",    TEAL,   TEALB),
        (7.94, "④ Retrieve",   GREEN,  GREENBG),
        (10.46,"⑤ Route",      PINK,   PINKBG),
    ]
    for x, label, fc, bg in headers:
        box(slide, x, 0.95, 2.38, 0.28, bg, fc, 0.8, radius=True)
        tb(slide, x + 0.08, 0.97, 2.22, 0.24, label, 9, fc, True, PP_ALIGN.CENTER)

    # ── Step 1: User Query ────────────────────────────────────────────────────
    box(slide, 0.38, 1.32, 2.38, 0.58, WHITE, BLUE, 1.2, radius=True)
    multi(slide, 0.42, 1.36, 2.30, 0.52, [
        ("Natural language question from the chat UI", 7.5, INK, False, PP_ALIGN.LEFT),
        ("session_id groups multi-turn history", 7, MUTED, False, PP_ALIGN.LEFT),
        ("trace_id generated at request start", 7, MUTED, False, PP_ALIGN.LEFT),
    ])

    # Example query pill
    box(slide, 0.38, 1.97, 2.38, 0.30, BLUEBG, BLUE, 0.6, radius=True)
    tb(slide, 0.46, 1.99, 2.22, 0.28,
       '"How do I submit a Spark job?"',
       7.5, BLUE, True, PP_ALIGN.CENTER)

    # ── Step 2: Classify ────────────────────────────────────────────────────────
    box(slide, 2.90, 1.32, 2.38, 0.95, WHITE, PURPLE, 1.2, radius=True)
    multi(slide, 2.94, 1.36, 2.30, 0.44, [
        ("LLM call — gpt-4o-mini, T=0.0", 7.5, INK, False, PP_ALIGN.LEFT),
        ("Returns: intent + confidence + reasoning", 7, MUTED, False, PP_ALIGN.LEFT),
    ])

    # 5 intent pills
    intent_data = [
        ("DOC_QA",          BLUE,   BLUEBG),
        ("ARTIFACT_SEARCH", GREEN,  GREENBG),
        ("USER_SEARCH",     AMBER,  AMBERBG),
        ("HYBRID",          PURPLE, PURPBG),
        ("OUT_OF_SCOPE",    RED,    RGBColor(254,226,226)),
    ]
    for i, (label, fc, bg) in enumerate(intent_data):
        box(slide, 2.93, 1.82 + i * 0.26, 2.32, 0.22, bg, fc, 0.6, radius=True)
        tb(slide, 2.97, 1.84 + i * 0.26, 2.24, 0.20, label, 7.5, fc, True, PP_ALIGN.CENTER)

    # out-of-scope short-circuit arrow label
    arr(slide, 5.28, 3.12, 10.46, 3.12, RED, 1.0)
    box(slide, 5.42, 3.00, 4.92, 0.22, RGBColor(254,226,226), RED, 0.6, radius=True)
    tb(slide, 5.46, 3.02, 4.86, 0.20,
       "OUT_OF_SCOPE → hardcoded reply, no further LLM calls, no scores",
       7, RED, True, PP_ALIGN.CENTER)

    # ── Step 3: Query Rewrite ────────────────────────────────────────────────────
    box(slide, 5.42, 1.32, 2.38, 0.95, WHITE, TEAL, 1.2, radius=True)
    multi(slide, 5.46, 1.36, 2.30, 0.88, [
        ("LLM call — gpt-4o-mini, T=0.0", 7.5, INK, False, PP_ALIGN.LEFT),
        ("Enriches query for better vector recall", 7, MUTED, False, PP_ALIGN.LEFT),
        ("max_tokens = 80", 7, MUTED, False, PP_ALIGN.LEFT),
        ('Input:  "Spark job Kubeflow"', 7, TEAL, False, PP_ALIGN.LEFT),
        ('Output: "submit Spark application\nKubeflow pipeline resource config"', 7, TEAL, False, PP_ALIGN.LEFT),
    ])

    # ── Step 4: Retriever details ────────────────────────────────────────────────
    retriever_data = [
        (7.94, 1.32, "platform_docs",       "chunk_text  ·  top-5",       "DOC_QA / HYBRID",    GREEN,  GREENBG),
        (7.94, 2.22, "artifact_summaries",  "artifact_summary  ·  top-5", "ARTIFACT / HYBRID",  AMBER,  AMBERBG),
        (7.94, 3.12, "user_profiles",       "user_profile  ·  top-5",     "USER / HYBRID",      PURPLE, PURPBG),
    ]
    for rx, ry, coll, fields, intents, fc, bg in retriever_data:
        box(slide, rx, ry, 2.38, 0.80, bg, fc, 0.8, radius=True)
        multi(slide, rx + 0.08, ry + 0.06, 2.22, 0.70, [
            (coll,   8.5, fc,    True,  PP_ALIGN.LEFT),
            (fields, 7,   DARKSL,False, PP_ALIGN.LEFT),
            (f"Used by: {intents}", 6.5, MUTED, False, PP_ALIGN.LEFT),
            ("HNSW · COSINE · 1536-dim", 6.5, MUTED, False, PP_ALIGN.LEFT),
        ])

    # Name-resolution shortcut box
    box(slide, 7.94, 4.02, 2.38, 0.58, PINKBG, PINK, 0.8, radius=True)
    multi(slide, 8.00, 4.06, 2.26, 0.52, [
        ("Name Resolver",  8.5, PINK,  True,  PP_ALIGN.LEFT),
        ("RapidFuzz string match", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("→ LLM disambiguation", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("→ Milvus profile fetch", 7, DARKSL, False, PP_ALIGN.LEFT),
    ])

    # ── Step 5: Routing summary ──────────────────────────────────────────────────
    route_data = [
        (1.32, "DOC_QA",          "doc_store → top 5",              BLUE,   BLUEBG),
        (2.00, "ARTIFACT_SEARCH", "artifact_store → top 5",         GREEN,  GREENBG),
        (2.68, "USER_SEARCH",     "user_store OR name resolver",     AMBER,  AMBERBG),
        (3.36, "HYBRID",          "all 3 stores → top 3 each",      PURPLE, PURPBG),
    ]
    for ry, label, detail, fc, bg in route_data:
        box(slide, 10.46, ry, 2.50, 0.58, bg, fc, 0.8, radius=True)
        multi(slide, 10.54, ry + 0.06, 2.34, 0.50, [
            (label,  8.5, fc,    True,  PP_ALIGN.LEFT),
            (detail, 7,   DARKSL,False, PP_ALIGN.LEFT),
        ])

    # ── Arrows: Query → Classify → Rewrite → Retrieve → Route ────────────────
    arr(slide, 2.76, 1.61, 2.90, 1.61, DARKSL, 1.2)
    arr(slide, 5.28, 1.61, 5.42, 1.61, DARKSL, 1.2)
    arr(slide, 7.80, 1.72, 7.94, 1.72, GREEN,  1.2)
    arr(slide, 7.80, 2.62, 7.94, 2.62, AMBER,  1.2)
    arr(slide, 7.80, 3.52, 7.94, 3.52, PURPLE, 1.2)
    arr(slide, 10.32, 1.61, 10.46, 1.61, DARKSL, 1.2)

    # Classify → Rewrite (with confidence note)
    box(slide, 3.80, 2.34, 1.54, 0.22, WHITE, SLATE, 0.5, radius=True)
    tb(slide, 3.84, 2.36, 1.46, 0.20, "confidence 0–1", 6.5, MUTED, False, PP_ALIGN.CENTER)

    # ── Legend / footnote ─────────────────────────────────────────────────────
    box(slide, 0.38, 4.72, 12.57, 0.30, LIGHT, SLATE, 0.5)
    multi(slide, 0.46, 4.74, 12.40, 0.28, [
        ("All LLM calls carry trace_id via LiteLLM extra_body → Langfuse groups every generation in one trace per request  ·  Embeddings: text-embedding-3-small · 1536-dim  ·  Vector search: COSINE similarity",
         7, MUTED, False, PP_ALIGN.LEFT),
    ])


# ── Slide 2: Generation Process ────────────────────────────────────────────────

def slide_generation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand(slide)

    # Title
    box(slide, 0.38, 0.42, 12.57, 0.46, GREENBG, GREEN, 0)
    tb(slide, 0.5, 0.45, 12.3, 0.40,
       "Generation Process — Context Assembly → LLM Generate → Observability Trace",
       12.5, INK, True, PP_ALIGN.LEFT)

    # ── Main pipeline row: 4 boxes ────────────────────────────────────────────
    stages = [
        (0.38, "Retrieved Context", GREEN,  GREENBG,
         ["doc chunks (chunk_text)", "artifact summaries", "user profiles",
          "Combined per intent type", "HYBRID gets all 3"]),
        (3.54, "Prompt Builder", BLUE,   BLUEBG,
         ["Selects template per intent:", "  doc_qa / artifact_search", "  user_search / hybrid",
          "Injects context + user query", "Prepends chat history"]),
        (6.70, "LLM Generate", PURPLE, PURPBG,
         ["gpt-4o-mini via LiteLLM proxy", "temperature = 0.2", "max_tokens = 600",
          "trace_id forwarded in extra_body", "Trace name: chat · {intent}"]),
        (9.86, "Response", TEAL,   TEALB,
         ["answer (str)", "intent + confidence", "trace_id (for feedback)",
          "raw_docs / raw_artifacts", "raw_users lists"]),
    ]
    for x, title, fc, bg, bullets in stages:
        box(slide, x, 0.98, 3.02, 2.28, bg, fc, 1.2, radius=True)
        tb(slide, x + 0.12, 1.02, 2.80, 0.30, title, 10, fc, True)
        yy = 1.38
        for b in bullets:
            dot(slide, x + 0.15, yy + 0.05, 0.08, fc)
            tb(slide, x + 0.28, yy, 2.66, 0.26, b, 7.5, DARKSL)
            yy += 0.26

    # Arrows between stages
    for x in [3.40, 6.56, 9.72]:
        arr(slide, x, 2.12, x + 0.14, 2.12, DARKSL, 1.4)

    # ── Prompt template detail (below the main row) ───────────────────────────
    box(slide, 0.38, 3.36, 12.57, 0.24, INK, INK, 0)
    tb(slide, 0.50, 3.37, 12.30, 0.22, "Prompt Templates per Intent", 8.5, WHITE, True)

    template_data = [
        (0.38,  "DOC_QA",          BLUE,   BLUEBG,
         "System: answer only from provided docs\nUser: {retrieved_chunks}\nQuestion: {user_query}"),
        (3.54,  "ARTIFACT_SEARCH", GREEN,  GREENBG,
         "System: list matching artifacts with IDs\nUser: {artifact_results}\nQuery: {user_query}"),
        (6.70,  "USER_SEARCH",     AMBER,  AMBERBG,
         "System: describe the most relevant person\nUser: {user_results}\nQuery: {user_query}"),
        (9.86,  "HYBRID",          PURPLE, PURPBG,
         "System: combine docs + artifacts + people\nUser: all 3 context blocks\nQuery: {user_query}"),
    ]
    for x, label, fc, bg, desc in template_data:
        box(slide, x, 3.66, 3.02, 1.18, bg, fc, 0.8, radius=True)
        tb(slide, x + 0.10, 3.70, 2.84, 0.24, label, 8.5, fc, True)
        tb(slide, x + 0.10, 3.96, 2.84, 0.82, desc, 7, DARKSL, False)

    # ── Observability trace detail ────────────────────────────────────────────
    box(slide, 0.38, 4.92, 12.57, 0.24, DARKSL, DARKSL, 0)
    tb(slide, 0.50, 4.93, 12.30, 0.22, "Langfuse Trace Structure (per request)", 8.5, WHITE, True)

    # Trace container
    box(slide, 0.38, 5.22, 12.57, 0.86, LIGHT, SLATE, 0.6, dash=True)

    trace_items = [
        (0.52,  "trace_id",      "UUID — groups all spans", BLUE,   BLUEBG),
        (2.90,  "[classify]",    "generation span · gpt-4o-mini · T=0", PURPLE, PURPBG),
        (5.28,  "[rewrite]",     "generation span · gpt-4o-mini · T=0", TEAL,   TEALB),
        (7.66,  "[generate]",    "generation span · T=0.2 · max 600 tok", GREEN,  GREENBG),
        (10.04, "trace metadata","intent · confidence · source_count · hits", AMBER,  AMBERBG),
    ]
    for x, label, detail, fc, bg in trace_items:
        box(slide, x, 5.30, 2.28, 0.68, bg, fc, 0.6, radius=True)
        tb(slide, x + 0.08, 5.34, 2.12, 0.24, label, 8, fc, True)
        tb(slide, x + 0.08, 5.56, 2.12, 0.36, detail, 6.5, DARKSL, False)

    # footnote
    box(slide, 0.38, 6.16, 12.57, 0.30, LIGHT, SLATE, 0.5)
    tb(slide, 0.46, 6.18, 12.40, 0.28,
       "LiteLLM proxy logs all token usage, cost, latency, prompt and response to Langfuse automatically via callback. "
       "No explicit Langfuse SDK calls in the engine. Session history injected between system prompt and final user message.",
       7, MUTED, False, PP_ALIGN.LEFT)


# ── Slide 3: Evaluation Process ────────────────────────────────────────────────

def slide_evaluation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    brand(slide)

    # Title
    box(slide, 0.38, 0.42, 12.57, 0.46, AMBERBG, AMBER, 0)
    tb(slide, 0.5, 0.45, 12.3, 0.40,
       "Evaluation Process — Layer 1 Heuristics · Layer 2 RAGAS · LLM Judge · User Feedback",
       12.5, INK, True, PP_ALIGN.LEFT)

    # ── Layer headers ─────────────────────────────────────────────────────────
    layer_headers = [
        (0.38,  3.54, "Layer 1 — Heuristic Scores",     BLUE,   "Inline · ~0 ms · always synchronous"),
        (6.92,  6.03, "Layer 2 — RAGAS + LLM Judge",    PURPLE, "Background thread · 10–30 s delay"),
        (10.08, 2.86, "User Feedback",                   PINK,   "User-initiated via frontend"),
    ]
    for x, w, label, fc, sub in layer_headers:
        box(slide, x, 0.97, w, 0.40, fc, fc, 0)
        tb(slide, x + 0.12, 0.99, w - 0.24, 0.22, label, 9.5, WHITE, True)
        tb(slide, x + 0.12, 1.18, w - 0.24, 0.18, sub, 7, RGBColor(200,220,255), False)

    # ── Layer 1 score cards ────────────────────────────────────────────────────
    l1_scores = [
        (0.38, "response_length",  BLUE,   BLUEBG,
         "Normalised answer length",
         "< 50 chars → 0.5   |   50–300 → ramp to 1.0\n300–2000 chars → 1.0   |   > 2000 → penalised"),
        (3.00, "has_content",      GREEN,  GREENBG,
         "Substantive answer check",
         "1.0 = no fallback phrase detected\n0.0 = 'couldn't find' or 'please try again'"),
        (0.38, "intent_confidence",PURPLE, PURPBG,
         "Classifier certainty",
         "Direct passthrough from classifier output\nLow value = boundary / ambiguous query"),
        (3.00, "source_count",     AMBER,  AMBERBG,
         "Retrieval coverage",
         "retrieved_sources / 5, capped at 1.0\n0.0 when no chunks were returned"),
    ]
    positions = [(0.38, 1.46), (3.00, 1.46), (0.38, 2.60), (3.00, 2.60)]
    for (x, y), (_, name, fc, bg, title, detail) in zip(positions, l1_scores):
        box(slide, x, y, 2.46, 1.04, bg, fc, 0.8, radius=True)
        tb(slide, x + 0.10, y + 0.06, 2.28, 0.24, name, 8.5, fc, True)
        tb(slide, x + 0.10, y + 0.30, 2.28, 0.20, title, 7.5, INK, False)
        tb(slide, x + 0.10, y + 0.50, 2.28, 0.48, detail, 7, MUTED, False)

    # NOT posted note
    box(slide, 0.38, 3.72, 5.46, 0.46, RGBColor(254,226,226), RED, 0.6, radius=True)
    multi(slide, 0.46, 3.76, 5.30, 0.42, [
        ("Not posted for: ", 7.5, RED, True, PP_ALIGN.LEFT),
        ("OUT_OF_SCOPE  (engine short-circuits before scoring)", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("Exact USER_SEARCH match  (early return before scoring block)", 7, DARKSL, False, PP_ALIGN.LEFT),
    ])

    # ── Layer 2 score cards ────────────────────────────────────────────────────
    l2_scores = [
        (6.92, 1.46, "faithfulness",     PURPLE, PURPBG,
         "Answer grounded in retrieved context",
         "RAGAS 0–1  |  eval model: gpt-4o\nLow = LLM hallucinated from training data"),
        (9.06, 1.46, "answer_relevancy", BLUE,   BLUEBG,
         "Answer addresses the question",
         "RAGAS 0–1  |  uses embedding similarity\nLow = off-topic or overly generic answer"),
        (6.92, 2.60, "context_precision",TEAL,   TEALB,
         "Top context is most relevant",
         "RAGAS 0–1  |  ranking quality metric\nLow = wrong chunks ranked highest"),
        (9.06, 2.60, "profile_relevance",PINK,   PINKBG,
         "Profile answers the question",
         "LLM judge  |  scores: 0 / 0.3 / 0.7 / 1.0\nOnly for exact USER_SEARCH name match"),
    ]
    for x, y, name, fc, bg, title, detail in l2_scores:
        box(slide, x, y, 2.98, 1.04, bg, fc, 0.8, radius=True)
        tb(slide, x + 0.10, y + 0.06, 2.80, 0.24, name, 8.5, fc, True)
        tb(slide, x + 0.10, y + 0.30, 2.80, 0.20, title, 7.5, INK, False)
        tb(slide, x + 0.10, y + 0.50, 2.80, 0.48, detail, 7, MUTED, False)

    # RAGAS applies to note
    box(slide, 6.92, 3.72, 5.10, 0.46, GREENBG, GREEN, 0.6, radius=True)
    multi(slide, 7.00, 3.76, 4.94, 0.42, [
        ("faithfulness / answer_relevancy / context_precision apply to: ", 7.5, GREEN, True, PP_ALIGN.LEFT),
        ("DOC_QA · ARTIFACT_SEARCH · semantic USER_SEARCH · HYBRID", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("profile_relevance applies only to: exact USER_SEARCH name match", 7, DARKSL, False, PP_ALIGN.LEFT),
    ])

    # ── User Feedback ─────────────────────────────────────────────────────────
    box(slide, 10.08, 1.46, 2.86, 1.04, PINKBG, PINK, 0.8, radius=True)
    tb(slide, 10.18, 1.50, 2.68, 0.24, "user_feedback", 8.5, PINK, True)
    tb(slide, 10.18, 1.74, 2.68, 0.20, "Thumbs up / down in UI", 7.5, INK, False)
    multi(slide, 10.18, 1.94, 2.68, 0.48, [
        ("👍  thumbs-up  → 1.0", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("👎  thumbs-down → 0.0", 7, DARKSL, False, PP_ALIGN.LEFT),
        ("POST /observability/feedback", 7, MUTED, False, PP_ALIGN.LEFT),
    ])

    # ── Score matrix ──────────────────────────────────────────────────────────
    box(slide, 0.38, 4.26, 12.57, 0.24, INK, INK, 0)
    tb(slide, 0.50, 4.27, 12.30, 0.22, "Score Availability Matrix per Intent Path", 8.5, WHITE, True)

    matrix_headers = ["Intent Path", "L1: heuristics", "L2: faithfulness", "L2: answer_rel.", "L2: ctx_precision", "L2: profile_rel.", "user_feedback"]
    col_xs = [0.38, 2.16, 4.12, 5.94, 7.76, 9.58, 11.40]
    col_ws = [1.74, 1.92, 1.78, 1.78, 1.78, 1.78, 1.55]

    # header row
    for i, (hx, hw, htext) in enumerate(zip(col_xs, col_ws, matrix_headers)):
        box(slide, hx, 4.56, hw - 0.04, 0.28, PANEL, SLATE, 0.5)
        tb(slide, hx + 0.06, 4.58, hw - 0.16, 0.26, htext, 7, INK, True, PP_ALIGN.CENTER)

    Y = "✓"
    N = "–"
    rows = [
        ("DOC_QA",              Y, Y, Y, Y, N, Y),
        ("ARTIFACT_SEARCH",     Y, Y, Y, Y, N, Y),
        ("USER_SEARCH (semantic)",Y,Y, Y, Y, N, Y),
        ("USER_SEARCH (exact)", N, N, N, N, Y, Y),
        ("HYBRID",              Y, Y, Y, Y, N, Y),
        ("OUT_OF_SCOPE",        N, N, N, N, N, N),
    ]
    row_colors = [BLUEBG, GREENBG, AMBERBG, PINKBG, PURPBG, RGBColor(254,226,226)]
    for ri, (row, rc) in enumerate(zip(rows, row_colors)):
        ry = 4.90 + ri * 0.30
        for ci, (cx, cw, cell) in enumerate(zip(col_xs, col_ws, row)):
            bg = rc if ci == 0 else WHITE
            fc_val = GREEN if cell == Y else (MUTED if cell == N else INK)
            box(slide, cx, ry, cw - 0.04, 0.26, bg, SLATE, 0.4)
            tb(slide, cx + 0.06, ry + 0.02, cw - 0.16, 0.24,
               cell, 7.5 if ci == 0 else 9, INK if ci == 0 else fc_val,
               ci == 0, PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_retrieval(prs)
    slide_generation(prs)
    slide_evaluation(prs)

    prs.save(OUTPUT)
    print(f"Saved → {OUTPUT}")


if __name__ == "__main__":
    main()
