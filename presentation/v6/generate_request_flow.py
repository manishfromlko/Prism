"""
Single slide: Chat Engine — full request flow diagram.

Paths captured (matching engine.py exactly):
  A  OUT_OF_SCOPE    → classify only → return hardcoded reply (no scores)
  B  USER_SEARCH     → classify → rewrite → RapidFuzz → LLM resolver
       B1 exact_uid  → fetch Milvus profile → LLM judge (profile_relevance) → return
       B2 ambiguous  → return disambiguation answer (no scores at all)
       B3 no-name    → falls through to vector retrieval → normal generate path
  C  Normal path     → classify → rewrite → vector retrieve → build prompt
                     → generate → Layer 1 scores → Layer 2 RAGAS (async) → return

Run:  python generate_request_flow.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Chat_Engine_Request_Flow.pptx"

WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(248, 250, 252)
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(100, 116, 139)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(180, 100, 0)
RED    = RGBColor(185, 28,  28)
PINK   = RGBColor(190, 24,  93)
PURPLE = RGBColor(109, 40,  217)
TEAL   = RGBColor(15,  118, 110)
SLATE  = RGBColor(148, 163, 184)
DARKSL = RGBColor(71,  85,  105)
PANEL  = RGBColor(241, 245, 249)
BLUEBG = RGBColor(219, 234, 254)
GREENBG= RGBColor(220, 252, 231)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)
TEALB  = RGBColor(204, 251, 241)
WARN   = RGBColor(254, 226, 226)
REDBRD = RGBColor(220, 38,  38)


# ── Primitives ─────────────────────────────────────────────────────────────────

def _shape(slide, kind, x, y, w, h):
    return slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))


def rnd(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=1.2):
    s = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(lw)
    s.adjustments[0] = 0.06
    return s


def dashed_rect(slide, x, y, w, h, fill=None, line=SLATE, lw=1.0):
    s = _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.color.rgb = line; s.line.width = Pt(lw)
    s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return s


def dmnd(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=1.4):
    """Diamond decision shape."""
    s = _shape(slide, MSO_SHAPE.DIAMOND, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(lw)
    return s


def oval(slide, x, y, w, h, fill=GREEN, line=GREEN, lw=1.2):
    s = _shape(slide, MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(lw)
    return s


def txt(slide, x, y, w, h, text, size=9, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color
    return b


def multi(slide, x, y, w, h, rows):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)
    return b


def arr(slide, x1, y1, x2, y2, color=DARKSL, lw=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(lw)
    return c


def label(slide, x, y, text, color=DARKSL, size=7.5, bold=False):
    txt(slide, x, y, 1.20, 0.22, text, size, color, bold, PP_ALIGN.CENTER)


def process_box(slide, x, y, w, h, title, sub, fill, border):
    """Box with bold title + smaller subtitle."""
    rnd(slide, x, y, w, h, fill, border, 1.4)
    multi(slide, x + 0.12, y + 0.04, w - 0.24, h - 0.08, [
        (title, 9.5, border, True),
        (sub,   8,   DARKSL, False),
    ])


def decision_box(slide, x, y, w, h, text, fill, border):
    """Diamond with text. Returns corner coords."""
    dmnd(slide, x, y, w, h, fill, border, 1.4)
    txt(slide, x + 0.10, y + 0.08, w - 0.20, h - 0.16,
        text, 8.5, border, True, PP_ALIGN.CENTER)
    cx = x + w / 2
    cy = y + h / 2
    return {
        "top":    (cx,     y),
        "bottom": (cx,     y + h),
        "left":   (x,      cy),
        "right":  (x + w,  cy),
    }


def terminal(slide, x, y, w, h, text, fill, border):
    """Rounded oval terminal (start/end)."""
    oval(slide, x, y, w, h, fill, border, 1.4)
    txt(slide, x + 0.08, y + 0.04, w - 0.16, h - 0.08,
        text, 9, WHITE, True, PP_ALIGN.CENTER)


def llm_badge(slide, x, y):
    """Small 'LLM' pill badge to mark an LLM call."""
    rnd(slide, x, y, 0.54, 0.20, PURPBG, PURPLE, 0.8)
    txt(slide, x + 0.04, y + 0.02, 0.46, 0.17, "LLM", 7, PURPLE, True, PP_ALIGN.CENTER)


# ── Main slide builder ─────────────────────────────────────────────────────────

def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── Title ──────────────────────────────────────────────────────────────────
    txt(slide, 0.28, 0.04, 12.8, 0.52,
        "Chat Engine  —  Request Flow",
        22, INK, True)
    txt(slide, 0.30, 0.56, 12.8, 0.24,
        "How a user query travels through classification → rewrite → retrieval → generation → scoring",
        9.5, MUTED, False)

    # ══════════════════════════════════════════════════════════════════════════
    # MAIN SPINE  (vertical, left column)
    # x=0.28, w=3.60   center x=2.08
    # ══════════════════════════════════════════════════════════════════════════
    SX = 0.28; SW = 3.60; SC = SX + SW / 2   # spine center x

    # [1] API Request
    terminal(slide, SX + 0.60, 0.88, SW - 1.20, 0.44, "User query received", BLUE, BLUE)
    arr(slide, SC, 1.32, SC, 1.46, DARKSL, 1.4)

    # [2] Generate trace_id
    process_box(slide, SX, 1.46, SW, 0.50,
                "Generate trace_id (UUID)",
                "Groups all LLM spans under one Langfuse trace",
                LIGHT, DARKSL)
    arr(slide, SC, 1.96, SC, 2.10, DARKSL, 1.4)

    # [3] Intent Classifier
    process_box(slide, SX, 2.10, SW, 0.54,
                "Intent Classifier",
                "LLM · gpt-4o-mini · T=0.0  →  intent, confidence, reasoning",
                PURPBG, PURPLE)
    llm_badge(slide, SX + SW - 0.62, 2.12)
    arr(slide, SC, 2.64, SC, 2.78, DARKSL, 1.4)

    # [D1] OUT_OF_SCOPE diamond
    d1 = decision_box(slide, SX, 2.78, SW, 0.56,
                      "OUT_OF_SCOPE ?", WARN, RED)
    # NO label (straight down)
    label(slide, SC - 0.64, 3.36, "No", DARKSL, 8, True)
    arr(slide, SC, 3.34, SC, 3.50, DARKSL, 1.4)

    # [4] Query Rewriter
    process_box(slide, SX, 3.50, SW, 0.54,
                "Query Rewriter",
                "LLM · gpt-4o-mini · T=0.0 · max 80 tok  →  enriched search_query",
                TEALB, TEAL)
    llm_badge(slide, SX + SW - 0.62, 3.52)
    arr(slide, SC, 4.04, SC, 4.18, DARKSL, 1.4)

    # [D2] USER_SEARCH diamond
    d2 = decision_box(slide, SX, 4.18, SW, 0.56,
                      "intent == USER_SEARCH ?", AMBERBG, AMBER)
    # NO label
    label(slide, SC - 0.64, 4.76, "No", DARKSL, 8, True)
    arr(slide, SC, 4.74, SC, 4.90, DARKSL, 1.4)

    # [5] Vector Retrieval
    rnd(slide, SX, 4.90, SW, 0.76, GREENBG, GREEN, 1.4)
    multi(slide, SX + 0.12, 4.94, SW - 0.24, 0.68, [
        ("Vector Retrieval  (Milvus similarity search)", 9.5, GREEN,  True),
        ("DOC_QA  →  platform_docs   top-5 chunks",     8,   DARKSL, False),
        ("ARTIFACT  →  artifact_summaries   top-5",     8,   DARKSL, False),
        ("USER_SEARCH  →  user_profiles   top-5",       8,   DARKSL, False),
        ("HYBRID  →  all 3 stores   top-3 each",        8,   DARKSL, False),
    ])
    arr(slide, SC, 5.66, SC, 5.80, DARKSL, 1.4)

    # [6] Build Prompt
    process_box(slide, SX, 5.80, SW, 0.50,
                "Build Prompt Messages",
                "Select template per intent · inject chat history",
                BLUEBG, BLUE)
    arr(slide, SC, 6.30, SC, 6.44, DARKSL, 1.4)

    # [7] LLM Generate
    process_box(slide, SX, 6.44, SW, 0.54,
                "LLM Generate",
                "gpt-4o-mini · T=0.2 · max 600 tokens · trace_name: chat · {intent}",
                PURPBG, PURPLE)
    llm_badge(slide, SX + SW - 0.62, 6.46)

    # ══════════════════════════════════════════════════════════════════════════
    # SCORING ROW  (below generate, spans full width)
    # ══════════════════════════════════════════════════════════════════════════
    arr(slide, SC, 6.98, SC, 7.10, DARKSL, 1.4)

    # Scoring strip at bottom
    rnd(slide, 0.28, 7.10, 5.50, 0.30, GREENBG, GREEN, 1.0)
    txt(slide, 0.38, 7.12, 5.30, 0.26,
        "Layer 1 scores  (inline · synchronous):  "
        "response_length · has_content · intent_confidence · source_count",
        8, GREEN, False)

    rnd(slide, 5.88, 7.10, 7.18, 0.30, PANEL, DARKSL, 0.8)
    txt(slide, 5.98, 7.12, 5.44, 0.26,
        "Layer 2 RAGAS  (background daemon thread · 10–30 s):  "
        "faithfulness · answer_relevancy · context_precision",
        8, DARKSL, False)
    txt(slide, 11.48, 7.12, 1.56, 0.26,
        "→ Return + trace_id", 8, BLUE, True)

    # ══════════════════════════════════════════════════════════════════════════
    # BRANCH A — OUT_OF_SCOPE  (top right)
    # ══════════════════════════════════════════════════════════════════════════
    BRX = 4.40   # branch right start x

    # Arrow from D1 right point → branch
    arr(slide, d1["right"][0], d1["right"][1], BRX, d1["right"][1], RED, 1.4)
    label(slide, d1["right"][0] + 0.24, d1["right"][1] - 0.24, "Yes", RED, 8, True)

    # OOS result box
    rnd(slide, BRX, 2.62, 8.62, 0.88, WARN, RED, 1.4)
    multi(slide, BRX + 0.14, 2.68, 8.34, 0.74, [
        ("Return hardcoded OUT_OF_SCOPE reply", 10, RED, True),
        ("• 1 LLM call only  (classify span)", 8.5, DARKSL, False),
        ("• No Query Rewriter, no retrieval", 8.5, DARKSL, False),
        ("• No scores posted to Langfuse  (engine returns before scoring block)", 8.5, DARKSL, False),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # BRANCH B — USER_SEARCH name resolution  (middle right)
    # ══════════════════════════════════════════════════════════════════════════
    NX = 4.40   # name resolution column x
    NW = 4.20   # column width
    NC = NX + NW / 2  # center

    # Arrow from D2 right point → name resolution column
    arr(slide, d2["right"][0], d2["right"][1], NX, d2["right"][1], AMBER, 1.4)
    label(slide, d2["right"][0] + 0.22, d2["right"][1] - 0.24, "Yes", AMBER, 8, True)

    # [N1] RapidFuzz name match
    process_box(slide, NX, 4.02, NW, 0.52,
                "RapidFuzz String Match",
                "retrieve_candidates(query, all_user_ids)  →  name_candidates list",
                AMBERBG, AMBER)
    arr(slide, NC, 4.54, NC, 4.68, AMBER, 1.4)

    # [D3] Candidates found?
    d3 = decision_box(slide, NX, 4.68, NW, 0.56,
                      "name_candidates found ?", AMBERBG, AMBER)
    label(slide, NC - 0.64, 5.26, "Yes", AMBER, 8, True)
    arr(slide, NC, 5.24, NC, 5.38, AMBER, 1.4)

    # "No" from D3 → back to main spine vector retrieve
    arr(slide, d3["left"][0], d3["left"][1], SX + SW, 5.28, DARKSL, 1.2)
    label(slide, NX - 0.94, 4.86, "No →\nfallback to\nvector search", DARKSL, 7, False)

    # [N2] LLM Name Resolver
    process_box(slide, NX, 5.38, NW, 0.52,
                "LLM Name Resolver  (user_resolve span)",
                "gpt-4o-mini · disambiguates candidates → resolved: {exact_uid, answer}",
                AMBERBG, AMBER)
    llm_badge(slide, NX + NW - 0.62, 5.40)
    arr(slide, NC, 5.90, NC, 6.04, AMBER, 1.4)

    # [D4] Exact UID?
    d4 = decision_box(slide, NX, 6.04, NW, 0.56,
                      "exact_uid resolved ?", AMBERBG, AMBER)

    # ── Path B1: exact match (right of D4) ────────────────────────────────────
    EX = 9.00; EW = 4.08
    arr(slide, d4["right"][0], d4["right"][1], EX, d4["right"][1], GREEN, 1.4)
    label(slide, d4["right"][0] + 0.22, d4["right"][1] - 0.24, "Yes", GREEN, 8, True)

    rnd(slide, EX, 6.04, EW, 0.52, GREENBG, GREEN, 1.4)
    multi(slide, EX + 0.12, 6.08, EW - 0.24, 0.44, [
        ("Fetch Milvus profile directly", 9.5, GREEN, True),
        ("user_store.get_profile(exact_uid)  →  raw profile text", 8, DARKSL, False),
    ])
    arr(slide, EX + EW / 2, 6.56, EX + EW / 2, 6.70, GREEN, 1.4)

    rnd(slide, EX, 6.70, EW, 0.52, GREENBG, GREEN, 1.4)
    multi(slide, EX + 0.12, 6.74, EW - 0.24, 0.44, [
        ("Layer 2: LLM Judge  →  profile_relevance", 9.5, GREEN, True),
        ("evaluate_in_background(exact_match=True) · no Layer 1 scores", 8, DARKSL, False),
    ])
    arr(slide, EX + EW / 2, 7.22, EX + EW / 2, 7.36, GREEN, 1.4)
    terminal(slide, EX + 0.54, 7.36, EW - 1.08, 0.26, "Return response", GREEN, GREEN)

    # ── Path B2: ambiguous (below D4) ─────────────────────────────────────────
    arr(slide, d4["bottom"][0], d4["bottom"][1],
        d4["bottom"][0], d4["bottom"][1] + 0.14, PINK, 1.4)
    label(slide, d4["bottom"][0] - 0.64, 6.62, "No", PINK, 8, True)

    rnd(slide, NX, 6.74, NW, 0.52, PINKBG, PINK, 1.4)
    multi(slide, NX + 0.12, 6.78, NW - 0.24, 0.44, [
        ("Return disambiguation answer", 9.5, PINK, True),
        ("confidence = 0.5  ·  NO Layer 1 scores  ·  NO Layer 2 scores", 8, DARKSL, False),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # LEGEND
    # ══════════════════════════════════════════════════════════════════════════
    lx = 9.00; ly = 0.86
    txt(slide, lx, ly, 4.10, 0.22, "Legend", 9, INK, True)

    legend_items = [
        (PURPBG, PURPLE, "LLM call  (classify / rewrite / generate / user_resolve)"),
        (GREENBG, GREEN,  "Scoring / retrieval step"),
        (AMBERBG, AMBER,  "USER_SEARCH name resolution path"),
        (WARN,    RED,    "OUT_OF_SCOPE path  (exits immediately)"),
        (PINKBG,  PINK,   "USER_SEARCH ambiguous  (exits, no scores)"),
    ]
    for i, (bg, fc, desc) in enumerate(legend_items):
        ly_i = 1.14 + i * 0.36
        rnd(slide, lx, ly_i, 0.36, 0.24, bg, fc, 1.0)
        txt(slide, lx + 0.44, ly_i + 0.02, 3.64, 0.22, desc, 8, DARKSL)

    # LLM badge in legend
    rnd(slide, lx, 1.14 + 5 * 0.36, 0.54, 0.22, PURPBG, PURPLE, 0.8)
    txt(slide, lx + 0.04, 1.14 + 5 * 0.36 + 0.02, 0.46, 0.18,
        "LLM", 7, PURPLE, True, PP_ALIGN.CENTER)
    txt(slide, lx + 0.62, 1.14 + 5 * 0.36 + 0.02, 3.40, 0.18,
        "Badge marks each LLM API call made", 8, DARKSL)

    # ══════════════════════════════════════════════════════════════════════════
    # Exit count summary
    # ══════════════════════════════════════════════════════════════════════════
    sx = 9.00; sy = 3.44
    dashed_rect(slide, sx, sy, 4.16, 1.44, PANEL, SLATE, 0.8)
    txt(slide, sx + 0.12, sy + 0.06, 3.90, 0.26, "Exit paths & LLM call count", 9, INK, True)
    exit_rows = [
        ("Path A  OUT_OF_SCOPE",       "1 call",  RED),
        ("Path B1  exact name match",  "3 calls", GREEN),
        ("Path B2  ambiguous name",    "3 calls", PINK),
        ("Path C   normal flow",       "3 calls", BLUE),
    ]
    for i, (path, calls, fc) in enumerate(exit_rows):
        ey = sy + 0.36 + i * 0.27
        txt(slide, sx + 0.12, ey, 2.80, 0.26, path,  8.5, fc,    True)
        txt(slide, sx + 2.98, ey, 1.06, 0.26, calls, 8.5, DARKSL, False)

    # Bottom accent bar
    for i, c in enumerate([BLUE, GREEN, AMBER, PINK]):
        rnd(slide, i * 3.3325, 7.68, 3.3325, 0.10, c, c, 0)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.80)   # slightly taller for the flow
    build(prs)
    prs.save(OUTPUT)
    print(f"Saved → {OUTPUT}")


if __name__ == "__main__":
    main()
