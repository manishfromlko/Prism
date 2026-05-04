"""
Architecture diagram — Kubeflow Workspace Intelligence.
Mirrors the "On-Prem LLM Serving and Observability" style:
  left = Observability Plane, right-top = Control Plane, right-bottom = Data Plane.

Run:  python generate_architecture_diagram.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Kubeflow_System_Architecture.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(250, 250, 252)
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(100, 116, 139)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(180, 100, 0)
RED    = RGBColor(185, 28,  28)
PINK   = RGBColor(190, 24,  93)
PURPLE = RGBColor(109, 40,  217)
TEAL   = RGBColor(15,  118, 110)
SLATE  = RGBColor(148, 163, 184)
DARKSL = RGBColor(71,  85,  105)
PANEL  = RGBColor(241, 245, 249)
BLUEBG = RGBColor(219, 234, 254)
GREENBG= RGBColor(220, 252, 231)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)
TEALB  = RGBColor(204, 251, 241)
WARN   = RGBColor(254, 226, 226)


# ── Primitives ────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=0.8, dash=False, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    if dash:
        s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius:
        s.adjustments[0] = 0.04
    return s


def txt(slide, x, y, w, h, text, size=10, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.margin_left   = Inches(0.04)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text      = text
    p.alignment = align
    p.font.size   = Pt(size)
    p.font.bold   = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return box


def multiline(slide, x, y, w, h, lines):
    """lines = [(text, size, color, bold), ...]"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left   = Inches(0.08)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size  = Pt(size)
        p.font.color.rgb = color
        p.font.bold  = bold
        p.alignment  = PP_ALIGN.LEFT
        p.space_after = Pt(1)
    return box


def labeled_box(slide, x, y, w, h, label, fill=WHITE, line=SLATE, lw=0.8,
                fsize=10.5, fcol=INK):
    s = rect(slide, x, y, w, h, fill, line, lw)
    tf = s.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left   = Inches(0.1)
    tf.margin_right  = Inches(0.08)
    tf.margin_top    = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size  = Pt(fsize)
    p.font.color.rgb = fcol
    p.font.bold = True
    return s


def arrow(slide, x1, y1, x2, y2, color=DARKSL, lw=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    return c


def bullet_list(slide, x, y, w, items, title=None, tcolor=INK, isize=9, icolor=MUTED,
                tcol_bold=True):
    yy = y
    if title:
        txt(slide, x, yy, w, 0.22, title, size=9, color=tcolor, bold=tcol_bold)
        yy += 0.23
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(x + 0.04), Inches(yy + 0.055), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = icolor
        dot.line.fill.background()
        txt(slide, x + 0.18, yy, w - 0.22, 0.22, item,
            size=isize, color=INK)
        yy += 0.23
    return yy


def plane_label(slide, x, y, w, h, label, color=RED):
    """Bottom-right label inside a dashed boundary box, matching reference style."""
    txt(slide, x + w - 2.2, y + h - 0.28, 2.1, 0.26,
        label, size=10, color=color, bold=True, align=PP_ALIGN.RIGHT)


def section_label(slide, x, y, w, label, color=DARKSL):
    """Thin bar label at top of a region."""
    b = rect(slide, x, y, w, 0.24, PANEL, PANEL, 0)
    b.line.fill.background()
    txt(slide, x + 0.1, y + 0.02, w - 0.15, 0.21,
        label, size=8, color=color, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
def make_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── Title ─────────────────────────────────────────────────────────────────
    txt(slide, 0.2, 0.06, 13.0, 0.52,
        "Kubeflow Workspace Intelligence — System Architecture",
        size=26, color=INK, bold=True)

    # ══════════════════════════════════════════════════════════════════════════
    # LEFT COLUMN — Observability Plane
    # ══════════════════════════════════════════════════════════════════════════
    OL, OT, OW, OH = 0.15, 0.63, 4.45, 5.92   # Observability Plane boundary

    obs_box = rect(slide, OL, OT, OW, OH, fill=None, line=SLATE, lw=1.5, dash=True,
                   radius=True)

    # — Langfuse ——————————————————————————————————————————————————————————————
    bullet_list(slide, OL + 0.22, OT + 0.2, OW - 0.3,
                ["request / response traces",
                 "scores (Layer 1 heuristics + RAGAS)",
                 "token usage & cost per call",
                 "intent, confidence, source count"],
                tcolor=MUTED)

    labeled_box(slide, OL + 0.95, OT + 1.23, 2.55, 0.62,
                "Langfuse", BLUEBG, BLUE, 1.2, fsize=12, fcol=BLUE)

    # callback arrow annotation
    txt(slide, OL + 3.6, OT + 1.42, 0.72, 0.28,
        "callback", size=8.5, color=MUTED, italic=True)

    # — Prometheus ————————————————————————————————————————————————————————————
    bullet_list(slide, OL + 0.22, OT + 2.1, OW - 0.3,
                ["API latency (p50 / p95 / p99)",
                 "request throughput & error rate",
                 "LiteLLM proxy queue depth"],
                tcolor=MUTED)

    labeled_box(slide, OL + 0.95, OT + 2.9, 2.55, 0.62,
                "Prometheus", GREENBG, GREEN, 1.2, fsize=12, fcol=GREEN)

    # — Arrow Prometheus → Grafana ————————————————————————————————————————————
    arrow(slide, OL + 2.22, OT + 3.52, OL + 2.22, OT + 3.75, DARKSL)

    # — Grafana ———————————————————————————————————————————————————————————————
    labeled_box(slide, OL + 0.95, OT + 3.75, 2.55, 0.62,
                "Grafana", AMBERBG, AMBER, 1.2, fsize=12, fcol=AMBER)

    bullet_list(slide, OL + 0.22, OT + 4.5, OW - 0.3,
                ["latency & cost dashboards",
                 "eval score trends (faithfulness,",
                 "  answer relevancy, context prec.)"],
                tcolor=MUTED)

    plane_label(slide, OL, OT, OW, OH, "Observability Plane")

    # ══════════════════════════════════════════════════════════════════════════
    # TOP RIGHT — Actor + Client (outside dashed boxes)
    # ══════════════════════════════════════════════════════════════════════════
    RL = 4.8   # right column left edge

    # Actor stick-figure (approximate with oval + rectangle + lines)
    ax, ay = 7.45, 0.64
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(ax), Inches(ay), Inches(0.28), Inches(0.28))
    head.fill.solid(); head.fill.fore_color.rgb = INK
    head.line.fill.background()
    arrow(slide, ax + 0.14, ay + 0.28, ax + 0.14, ay + 0.72, INK, 1.8)
    # arms
    arrow(slide, ax - 0.18, ay + 0.40, ax + 0.46, ay + 0.40, INK, 1.4)
    # legs
    arrow(slide, ax + 0.14, ay + 0.72, ax - 0.1, ay + 1.0, INK, 1.4)
    arrow(slide, ax + 0.14, ay + 0.72, ax + 0.38, ay + 1.0, INK, 1.4)
    txt(slide, ax - 0.18, ay + 1.02, 0.62, 0.22, "Actor",
        size=8.5, color=INK, align=PP_ALIGN.CENTER)

    # Arrow Actor → Next.js
    arrow(slide, ax + 0.14, ay + 1.22, ax + 0.14, ay + 1.55, INK, 1.5)

    # Client (sdk / curl) dashed box
    client_x = 10.55
    cb = rect(slide, client_x, 0.68, 2.55, 0.42, fill=None, line=SLATE, lw=1.0, dash=True)
    txt(slide, client_x + 0.1, 0.71, 2.35, 0.35,
        "Client  ( sdk, curl, Next.js BFF )",
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Arrow Client → (points down-left to Next.js)
    arrow(slide, client_x + 1.28, 1.10, client_x + 1.28, 1.55, MUTED, 1.2)

    # ══════════════════════════════════════════════════════════════════════════
    # RIGHT-TOP — Control Plane
    # ══════════════════════════════════════════════════════════════════════════
    CPL, CPT, CPW, CPH = RL, 1.6, 8.38, 3.35

    rect(slide, CPL, CPT, CPW, CPH, fill=None, line=SLATE, lw=1.5, dash=True, radius=True)

    # ── Next.js BFF ────────────────────────────────────────────────────────
    labeled_box(slide, CPL + 1.5, CPT + 0.22, 2.55, 0.62,
                "Next.js  BFF", BLUEBG, BLUE, 1.2, fsize=12, fcol=BLUE)

    # TLS / Rate Limit annotation (dashed, top right)
    tlsb = rect(slide, CPL + 5.05, CPT + 0.22, 2.88, 0.42,
                fill=None, line=SLATE, lw=0.9, dash=True)
    txt(slide, CPL + 5.1, CPT + 0.26, 2.78, 0.35,
        "( TLS / Global Rate Limit )",
        size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Arrow Next.js → FastAPI
    arrow(slide, CPL + 2.78, CPT + 0.84, CPL + 2.78, CPT + 1.14, INK, 1.5)

    # ── FastAPI + ChatEngine ───────────────────────────────────────────────
    labeled_box(slide, CPL + 1.5, CPT + 1.14, 2.55, 0.68,
                "FastAPI\nRetrieval API", GREENBG, GREEN, 1.2, fsize=10.5, fcol=GREEN)

    # Arrow FastAPI → LiteLLM (down to AI Gateway region)
    arrow(slide, CPL + 2.78, CPT + 1.82, CPL + 2.78, CPT + 2.12, INK, 1.5)

    # ── LiteLLM ───────────────────────────────────────────────────────────
    labeled_box(slide, CPL + 1.5, CPT + 2.12, 2.55, 0.68,
                "LiteLLM", PURPBG, PURPLE, 1.2, fsize=12, fcol=PURPLE)

    # AI Gateway annotation box (right of LiteLLM)
    agb = rect(slide, CPL + 4.3, CPT + 1.5, 3.75, 1.75,
               fill=None, line=SLATE, lw=0.9, dash=True)
    txt(slide, CPL + 4.4, CPT + 1.52, 3.55, 0.25,
        "AI Gateway", size=9.5, color=PURPLE, bold=True)
    bullet_list(slide, CPL + 4.38, CPT + 1.77, 3.55,
                ["Auth & API key management",
                 "Model routing & fallbacks",
                 "Observability callbacks",
                 "Model-aware rate limiting",
                 "Request / response logging"],
                isize=8.8, icolor=PURPLE)

    # LiteLLM → AI Gateway annotation arrow
    arrow(slide, CPL + 4.05, CPT + 2.46, CPL + 4.3, CPT + 2.46, PURPLE, 1.2)

    # ChatEngine / Intent Router annotation (below FastAPI)
    ceb = rect(slide, CPL + 4.3, CPT + 0.2, 3.72, 1.15,
               fill=None, line=SLATE, lw=0.9, dash=True)
    txt(slide, CPL + 4.4, CPT + 0.22, 3.52, 0.25,
        "ChatEngine Orchestrator", size=9.5, color=GREEN, bold=True)
    bullet_list(slide, CPL + 4.38, CPT + 0.47, 3.5,
                ["Intent classifier (DOC_QA / ARTIFACT /",
                 "  USER_SEARCH / HYBRID / OUT_OF_SCOPE)",
                 "Query rewriter + retriever router",
                 "Prompt builder + answer formatter"],
                isize=8.8, icolor=GREEN)

    # FastAPI → ChatEngine annotation arrow
    arrow(slide, CPL + 4.05, CPT + 1.48, CPL + 4.3, CPT + 0.9, GREEN, 1.2)

    plane_label(slide, CPL, CPT, CPW, CPH, "Control Plane")

    # ── Callback arrow: LiteLLM → Langfuse (crosses planes) ──────────────
    lx = CPL + 1.5   # LiteLLM left edge
    ly = CPT + 2.46  # LiteLLM mid-y
    # horizontal line leftward from LiteLLM to the observability plane
    arrow(slide, lx, ly, OL + OW, OT + 1.54, PURPLE, 1.3)

    # ══════════════════════════════════════════════════════════════════════════
    # RIGHT-BOTTOM — Data Plane  (Milvus + OpenAI)
    # ══════════════════════════════════════════════════════════════════════════
    DPL, DPT, DPW, DPH = RL, 5.1, 8.38, 2.2

    rect(slide, DPL, DPT, DPW, DPH, fill=None, line=SLATE, lw=1.5, dash=True, radius=True)

    # Arrow from LiteLLM down into Data Plane
    arrow(slide, CPL + 2.78, CPT + CPH, CPL + 2.78, DPT + 0.2, INK, 1.5)

    # ── 4 Milvus Collections (2x2 grid) ───────────────────────────────────
    cx, cy, cw, ch = DPL + 0.22, DPT + 0.28, 2.55, 0.75

    # Collection configs: (col, row, label, sub, fill, line, fc)
    colls = [
        (0, 0, "kubeflow_artifacts", "Artifact semantic search\n(raw chunks, HNSW index)", AMBERBG, AMBER, AMBER),
        (1, 0, "artifact_summaries",  "LLM-summarised notebooks\n& scripts (gpt-4o-mini)",     AMBERBG, AMBER, AMBER),
        (0, 1, "user_profiles",       "Workspace owner expertise\nprofiles (gpt-4o-mini)",      AMBERBG, AMBER, AMBER),
        (1, 1, "platform_docs",       "Kubeflow onboarding &\nhow-to doc chunks (.docx)",       AMBERBG, AMBER, AMBER),
    ]

    for col, row, name, sub, fill, line, fc in colls:
        bx = cx + col * (cw + 0.22)
        by = cy + row * (ch + 0.18)
        b = rect(slide, bx, by, cw, ch, fill, line, 1.2)
        multiline(slide, bx + 0.08, by + 0.06, cw - 0.16, ch - 0.1,
                  [(name, 9.5, fc, True),
                   (sub,  8.5, MUTED, False)])

    # Milvus label
    mw = 2 * cw + 0.22
    txt(slide, cx + mw / 2 - 0.6, DPT + DPH - 0.28, 1.2, 0.24,
        "Milvus", size=9.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    # ── OpenAI API box (right side) ───────────────────────────────────────
    oax = DPL + 5.7
    oa = rect(slide, oax, DPT + 0.28, 2.42, 1.68, PURPBG, PURPLE, 1.2, radius=True)
    multiline(slide, oax + 0.12, DPT + 0.34, 2.18, 1.58,
              [("OpenAI API", 11, PURPLE, True),
               ("gpt-4o-mini", 9, INK, True),
               ("  chat, summaries, profiles", 8.5, MUTED, False),
               ("text-embedding-3-small", 9, INK, True),
               ("  all vector embeddings", 8.5, MUTED, False),
               ("gpt-4o  (RAGAS eval judge)", 9, INK, True)])

    # Arrow LiteLLM → OpenAI (down-right)
    arrow(slide, CPL + 3.95, CPT + 2.46, oax + 1.21, DPT + 0.28, PURPLE, 1.2)

    plane_label(slide, DPL, DPT, DPW, DPH, "LLM Serving :  Data Plane")

    # ══════════════════════════════════════════════════════════════════════════
    # BOTTOM-LEFT — Supporting Infrastructure
    # ══════════════════════════════════════════════════════════════════════════
    txt(slide, OL + 0.05, OT + OH + 0.15, OW - 0.1, 0.22,
        "Supporting Infrastructure",
        size=8.5, color=INK, bold=False)
    infra = [
        "- Docker Compose (Milvus, LiteLLM + Postgres, Langfuse stack)",
        "- Python 3.11 / uvicorn / FastAPI — port 8000",
        "- Node.js 18+ / Next.js 15 — port 3002",
        "- Airflow DAG: Ingest → Embed → Summarise → Profile",
    ]
    yy = OT + OH + 0.37
    for line in infra:
        txt(slide, OL + 0.05, yy, OW - 0.1, 0.22, line,
            size=8.5, color=MUTED)
        yy += 0.22

    # ══════════════════════════════════════════════════════════════════════════
    # LAYER 2 EVAL annotation (connects RAGAS → Langfuse)
    # ══════════════════════════════════════════════════════════════════════════
    # Small RAGAS box inside observability plane
    rb = rect(slide, OL + 0.3, OT + 5.2, 3.4, 0.56, PINKBG, PINK, 1.0, radius=True)
    multiline(slide, OL + 0.38, OT + 5.22, 3.22, 0.5,
              [("Layer 2 — RAGAS + LLM-as-Judge (async background)", 8.5, PINK, True),
               ("faithfulness  ·  answer_relevancy  ·  context_precision  ·  profile_relevance",
                7.5, MUTED, False)])

    # Arrow from RAGAS box → Grafana label area (scores flow up to Langfuse)
    arrow(slide, OL + 2.0, OT + 5.2, OL + 2.22, OT + 4.37, PINK, 1.2)

    slide.notes_slide.notes_text_frame.text = (
        "Architecture mirrors the On-Prem LLM Serving diagram:\n"
        "LEFT = Observability Plane (Langfuse for traces/scores, Prometheus+Grafana for ops metrics, RAGAS for eval)\n"
        "RIGHT-TOP = Control Plane (Next.js BFF → FastAPI → ChatEngine intent router → LiteLLM AI Gateway)\n"
        "RIGHT-BOTTOM = Data Plane (4 Milvus collections in 2x2 + OpenAI API external)\n"
        "BOTTOM-LEFT = Supporting Infrastructure (Docker, Airflow ingestion DAG)\n"
        "Key data flow: User → Next.js → FastAPI → LiteLLM → OpenAI/Milvus → response\n"
        "Observability: LiteLLM sends callback to Langfuse; RAGAS runs async post-response."
    )


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    make_slide(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
