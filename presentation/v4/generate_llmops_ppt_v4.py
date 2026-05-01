from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = "Kubeflow_Workspace_Intelligence_LLMOps_Deck_v4.pptx"

INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
PINK = RGBColor(219, 39, 119)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


slides = [
    (
        "Kubeflow Workspace Intelligence",
        [
            "RAG assistant for workspace knowledge, notebook discovery, and user expertise profiling",
            "Built for data scientists and ML engineers",
            "FastAPI, Next.js, Milvus, LiteLLM, Langfuse, RAGAS",
        ],
        "Title slide. Frame this as a production-style LLMOps system, not a generic chatbot wrapper.",
    ),
    (
        "Problem",
        [
            "Kubeflow workspaces become knowledge silos as notebooks and scripts grow",
            "Filename search misses ML intent: tools, models, datasets, techniques",
            "Teams need grounded answers, expertise discovery, and operational visibility",
        ],
        "The pain is discovery plus trust. Notebooks hold institutional memory, but that memory is hard to query and hard to evaluate.",
    ),
    (
        "Solution Overview",
        [
            "Intent-routed RAG assistant over docs, artifacts, and user profiles",
            "Semantic retrieval backed by OpenAI embeddings and Milvus collections",
            "LLMOps loop with LiteLLM traces, Langfuse scores, and RAGAS evaluation",
        ],
        "The key architectural move is separate retrieval surfaces rather than one mixed index.",
    ),
    (
        "Dataset and Ingestion Scope",
        [
            "Current catalog: 5 workspaces, 211 artifacts",
            "132 notebooks, 75 scripts, 4 text files",
            "Full and incremental ingestion with content-hash change detection",
        ],
        "The prompt mentions approximately 250 notebooks, but the checked-in catalog has 132. Present the concrete number and note that the design scales.",
    ),
    (
        "Core Capabilities",
        [
            "Workspace browsing and profiling",
            "Notebook/script semantic search",
            "LLM-generated artifact summaries and user profiles",
            "Platform documentation QA and enterprise assistant",
        ],
        "The codebase includes platform-doc QA and assistant routing in addition to the prompt's original capabilities.",
    ),
    (
        "High-Level Architecture",
        [
            "Next.js UI uses server-side API routes as a BFF",
            "FastAPI owns retrieval, chat orchestration, admin sync, and feedback endpoints",
            "Milvus stores artifact chunks, summaries, profiles, and platform docs",
        ],
        "The frontend does not couple browser calls directly to Python. FastAPI is the orchestration boundary.",
    ),
    (
        "Ingestion and Indexing Pipeline",
        [
            "Scan workspace directories and classify supported files",
            "Extract notebook/script metadata and cell text",
            "Chunk by content type, embed, and index into Milvus",
        ],
        "Chunking is content-aware: recursive for notebooks, code splitter for scripts, markdown splitter for markdown.",
    ),
    (
        "Retrieval Surfaces",
        [
            "kubeflow_artifacts: raw artifact chunks",
            "artifact_summaries: concise notebook/script descriptions",
            "user_profiles: expertise summaries",
            "platform_docs: Word document chunks",
        ],
        "Separate stores improve precision and prompt quality at the cost of more indexing workflows.",
    ),
    (
        "Chatbot Runtime Flow",
        [
            "Create trace_id, classify intent, rewrite query",
            "Retrieve from intent-specific stores",
            "Build prompt, generate through LiteLLM, format structured response",
        ],
        "The same trace ID is propagated across classify, rewrite, and generate so Langfuse shows one coherent request.",
    ),
    (
        "Intent Routing",
        [
            "DOC_QA, ARTIFACT_SEARCH, USER_SEARCH, HYBRID, OUT_OF_SCOPE",
            "USER_SEARCH tries name resolution before semantic retrieval",
            "OUT_OF_SCOPE avoids unsupported real-time or external answers",
        ],
        "This is a practical hallucination-control and cost-control mechanism.",
    ),
    (
        "LLM Orchestration",
        [
            "All chat calls use an OpenAI-compatible client pointed at LiteLLM",
            "Generation defaults to gpt-4o-mini",
            "Evaluation and judge paths use gpt-4o",
        ],
        "LiteLLM centralizes model routing, provider config, token/cost visibility, and Langfuse callbacks.",
    ),
    (
        "Layer 1 Observability",
        [
            "LiteLLM forwards prompts, responses, latency, cost, tokens, and errors to Langfuse",
            "Application posts response_length, has_content, intent_confidence, source_count",
            "Trace ID is returned for feedback and debugging",
        ],
        "Layer 1 tells us what happened operationally and what the request cost.",
    ),
    (
        "Layer 2 RAG Evaluation",
        [
            "Runs after response in a background thread",
            "RAGAS: faithfulness, answer_relevancy, context_precision",
            "Exact USER_SEARCH uses LLM judge profile_relevance",
        ],
        "Layer 2 tells us whether the answer was grounded, relevant, and supported by useful context.",
    ),
    (
        "Evaluation Interpretation",
        [
            "Low faithfulness: unsupported claims",
            "Low answer_relevancy: answer missed the question",
            "Low context_precision: retrieval returned weak context",
        ],
        "These metrics are diagnostic. They tell you whether to tune prompts, retrieval, or routing.",
    ),
    (
        "API and Web Surface",
        [
            "/query, /chat, /workspaces, /user-profiles, /profile/workspace/{id}",
            "Admin endpoints for indexing, summaries, profiles, and doc ingestion",
            "Feedback and score endpoints attach quality signals to Langfuse traces",
        ],
        "The admin endpoints show lifecycle thinking: the system can rebuild and refresh its knowledge stores.",
    ),
    (
        "Design Decisions",
        [
            "Separate vector collections over one mixed index",
            "Summary-first profiles over raw-code profile generation",
            "Async evaluation over inline RAGAS",
            "LLM classifier over fixed keyword routing",
        ],
        "Each decision has a tradeoff: precision vs operational complexity, latency vs immediacy, flexibility vs model dependency.",
    ),
    (
        "Failure Cases",
        [
            "Missing/stale indexes, weak retrieval, wrong intent, hallucinated answer",
            "Ambiguous people queries and LLM/provider failures",
            "Mitigations: sync endpoints, query rewrite, OUT_OF_SCOPE, exact user resolution, traces, scores",
        ],
        "A strong interview answer names failure modes and ties each one to a detection or mitigation mechanism.",
    ),
    (
        "Production Readiness",
        [
            "Implemented: ingestion, indexing, assistant, traces, scores, RAGAS, feedback, admin sync",
            "Next: auth/ACLs, scheduled eval sets, queue-backed eval workers, SLO dashboards",
            "Prompt/model versioning should become part of release management",
        ],
        "Be honest: production-style now, production-final after identity, ACLs, quality gates, and deployment hardening.",
    ),
    (
        "Demo Walkthrough",
        [
            "Show workspaces and profile",
            "Run semantic artifact search",
            "Ask DOC_QA, ARTIFACT_SEARCH, and USER_SEARCH questions",
            "Open Langfuse trace and show scores",
        ],
        "Demo the loop: user question, routed answer, trace, quality score, diagnosis.",
    ),
    (
        "Close",
        [
            "Turns workspace artifacts into a queryable intelligence layer",
            "Separates ingestion, retrieval, orchestration, observability, and evaluation",
            "Built to explain failures, not just return answers",
        ],
        "The strongest takeaway is engineering discipline: observable, evaluable, and grounded RAG.",
    ),
]


def add_textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_brand(slide, index):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    for i, color in enumerate([BLUE, GREEN, AMBER, PINK]):
        accent = slide.shapes.add_shape(1, Inches(i * 3.33), Inches(0.32), Inches(3.33), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
    add_textbox(
        slide,
        0.65,
        7.08,
        12.0,
        0.28,
        f"Kubeflow Workspace Intelligence LLMOps - v4 - {index:02d}",
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


def make_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_brand(slide, 1)
    add_textbox(slide, 0.8, 1.6, 11.8, 0.85, "Kubeflow Workspace Intelligence", 38, INK, True)
    add_textbox(slide, 0.82, 2.55, 11.2, 0.55, "RAG-based chatbot with LLMOps observability and evaluation", 22, BLUE, False)
    add_textbox(
        slide,
        0.82,
        3.35,
        10.7,
        1.1,
        "Production-style assistant for workspace knowledge, notebook discovery, platform docs QA, and user expertise profiling.",
        20,
        MUTED,
    )
    add_textbox(slide, 0.82, 5.45, 11.5, 0.4, "FastAPI | Next.js | Milvus | LiteLLM | Langfuse | RAGAS", 17, INK, True)
    add_notes(slide, slides[0][2])


def make_bullet(prs, index, title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide, index)
    add_textbox(slide, 0.65, 0.72, 11.9, 0.55, title, 30, INK, True)

    y = 1.65
    for bullet in bullets:
        dot = slide.shapes.add_shape(9, Inches(0.85), Inches(y + 0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_textbox(slide, 1.12, y, 11.3, 0.52, bullet, 19, INK)
        y += 0.72
    add_notes(slide, notes)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    make_title(prs)
    for i, (title, bullets, notes) in enumerate(slides[1:], start=2):
        make_bullet(prs, i, title, bullets, notes)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
